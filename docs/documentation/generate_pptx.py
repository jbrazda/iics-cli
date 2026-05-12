"""
Generate iics-cli presentation as a PowerPoint file.
Run: python3 generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ───────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x00, 0x33, 0x66)   # headings / title bg
MID_BLUE    = RGBColor(0x00, 0x66, 0xCC)   # accents, table header
ORANGE      = RGBColor(0xFF, 0x66, 0x00)   # highlight colour
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MID_GRAY    = RGBColor(0x88, 0x88, 0x88)
GREEN       = RGBColor(0x00, 0x99, 0x44)
RED         = RGBColor(0xCC, 0x00, 0x00)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)   # dark code background
CODE_FG     = RGBColor(0xA6, 0xE3, 0xA1)   # green code text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, italic=False,
                 color=DARK_GRAY, align=PP_ALIGN.LEFT,
                 word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_multiline_text(slide, lines, x, y, w, h,
                       font_size=16, color=DARK_GRAY,
                       font_name="Calibri", line_spacing=None):
    """lines: list of (text, bold, italic, color_override)"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, italic, col = item, False, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            italic = item[2] if len(item) > 2 else False
            col = item[3] if len(item) > 3 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if line_spacing:
            p.space_before = Pt(line_spacing)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = col
        run.font.name = font_name
    return txBox

def slide_header(slide, title, subtitle=None):
    """Blue header bar across the top of the slide."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
    add_text_box(slide, title,
                 Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.7),
                 font_size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.3), Inches(0.75), Inches(12.5), Inches(0.4),
                     font_size=14, color=RGBColor(0xBB, 0xCC, 0xFF),
                     align=PP_ALIGN.LEFT)

def footer(slide, text="iics-cli  |  Informatica IICS Command Line Interface"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), DARK_BLUE)
    add_text_box(slide, text,
                 Inches(0.3), Inches(7.18), Inches(12.7), Inches(0.3),
                 font_size=9, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.LEFT)

def code_block(slide, code, x, y, w, h, font_size=11):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_text_box(slide, code, x + Inches(0.15), y + Inches(0.1),
                 w - Inches(0.3), h - Inches(0.2),
                 font_size=font_size, color=CODE_FG,
                 font_name="Courier New")

def bullet_list(slide, items, x, y, w, h, font_size=16,
                indent=Inches(0.3), color=DARK_GRAY,
                bullet="•"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col = item[2] if len(item) > 2 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{bullet}  {text}"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = "Calibri"
    return txBox


def add_table(slide, headers, rows, x, y, w, h,
              header_color=MID_BLUE, alt_row=True):
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    tbl = slide.shapes.add_table(num_rows, num_cols, x, y, w, h).table

    col_w = w // num_cols
    for i in range(num_cols):
        tbl.columns[i].width = col_w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

    # Data rows
    for ri, row in enumerate(rows):
        bg = LIGHT_GRAY if (alt_row and ri % 2 == 0) else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            # colour Yes/No
            if val == "Yes":
                run.font.color.rgb = GREEN
                run.font.bold = True
            elif val == "No":
                run.font.color.rgb = RED
            else:
                run.font.color.rgb = DARK_GRAY
            run.text = val
            run.font.size = Pt(12)
            run.font.name = "Calibri"
    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
add_rect(s, 0, Inches(3.5), SLIDE_W, Inches(0.06), ORANGE)

add_text_box(s, "iics-cli",
             Inches(1), Inches(1.2), Inches(11), Inches(1.4),
             font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "A Modern Command-Line Interface for Informatica IICS",
             Inches(1), Inches(2.6), Inches(11), Inches(0.8),
             font_size=24, color=RGBColor(0xBB, 0xCC, 0xFF),
             align=PP_ALIGN.CENTER)
add_text_box(s, "REST API v3  |  Go Binary  |  POSIX Compliant  |  Actively Maintained",
             Inches(1), Inches(3.7), Inches(11), Inches(0.6),
             font_size=18, color=RGBColor(0xFF, 0xAA, 0x44),
             align=PP_ALIGN.CENTER)
add_text_box(s, "30+ resource types  |  OS Keychain  |  Package management  |  CI/CD release workflows",
             Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             font_size=16, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – Agenda
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Agenda")
footer(s)

items_l = [
    "1.   The Problem – Limitations of the Official Informatica CLI",
    "2.   POSIX Exit Codes – Reliable CI/CD Integration",
    "3.   Machine-Readable Output – JSON, CSV, YAML, Tables & Themes",
    "4.   Secure Credential Management + OS Keychain",
    "5.   Named Environment Profiles + Interactive Wizard",
    "6.   Session Caching",
]
items_r = [
    "7.   Broad API Coverage – 27+ Resources",
    "8.   Package Command – Expand, Create, Inspect",
    "9.   Release Workflow – CI/CD Manifests & Plans",
    "10.  Publish & Unpublish CAI Assets",
    "11.  Full Region Support",
    "12.  Feature Comparison Summary",
]

bullet_list(s, items_l, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5),
            font_size=17, bullet="")
bullet_list(s, items_r, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5),
            font_size=17, bullet="")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "The Problem", "The Official Informatica CLI – what it cannot do")
footer(s)

problems = [
    ("No longer maintained – last update years ago, no bug fixes", True, RED),
    ("Covers only 5 operations: export, import, publish, list, version", False, DARK_GRAY),
    ("Always exits with code 0 – CI/CD pipelines cannot detect failures", False, DARK_GRAY),
    ("No structured output – cannot pipe to jq, grep, or other UNIX tools", False, DARK_GRAY),
    ("Credentials passed as command-line flags (-u user -p pass) – visible in ps, logs, history", False, DARK_GRAY),
    ("No environment profiles – must retype credentials for every org/env switch", False, DARK_GRAY),
    ("Authenticates on every invocation – no session reuse", False, DARK_GRAY),
    ("Only 3 hardcoded regions: us / eu / ap – many IICS pods not covered", False, DARK_GRAY),
    ("Limited to export/import/publish operations only – no user, role, connection, or schedule management", False, DARK_GRAY),
    ("Flat command model – no subcommand hierarchy, poor discoverability", False, DARK_GRAY),
]
bullet_list(s, problems, Inches(0.5), Inches(1.25), Inches(12.5), Inches(5.8),
            font_size=15)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – POSIX Exit Codes
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "POSIX Exit Codes", "Reliable failure detection in scripts and CI/CD pipelines")
footer(s)

add_text_box(s, "Official CLI – exit code is always 0",
             Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=RED)
code_block(s, (
    "# Export fails – but exit code is still 0\n"
    "iics export -u user@co.com -p secret -r us \\\n"
    "  -a \"MyProject/BadAsset.Process\"\n"
    "echo $?   # -> 0  (pipeline keeps running!)"
), Inches(0.5), Inches(1.65), Inches(5.8), Inches(1.6), font_size=11)

add_text_box(s, "iics-cli – three distinct exit codes",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
code_block(s, (
    "iics export create --name nightly --project Prod\n"
    "echo $?   # -> 1  (API / runtime error)\n\n"
    "iics export create --bad-flag\n"
    "echo $?   # -> 2  (usage / flag error)"
), Inches(7.0), Inches(1.65), Inches(5.8), Inches(1.6), font_size=11)

add_table(s,
    ["Exit Code", "Meaning", "Use Case"],
    [
        ["0", "Success", "All pipeline steps continue"],
        ["1", "Runtime / API error", "API call failed, network error"],
        ["2", "Usage error", "Unknown flag, missing argument"],
    ],
    Inches(2.0), Inches(3.5), Inches(9.0), Inches(1.7))

add_text_box(s,
    "Works correctly with set -e scripts, make targets, GitHub Actions, and Jenkins.",
    Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
    font_size=14, italic=True, color=MID_BLUE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – Output Formats
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Machine-Readable Output", "JSON, CSV, YAML, and Table formats – pipe into any UNIX tool")
footer(s)

add_text_box(s, "Every command supports --output / -o  table | json | csv | yaml",
             Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.4),
             font_size=15, bold=True, color=DARK_BLUE)

code_block(s, (
    "# Default table – human readable\n"
    "iics user list\n\n"
    "# JSON – pipe into jq\n"
    "iics user list -o json | jq -r '.[].userName'\n\n"
    "# CSV – import into Excel / Google Sheets\n"
    "iics connection list -o csv > connections.csv\n\n"
    "# YAML – human-friendly structured output\n"
    "iics schedule list -o yaml\n\n"
    "# Find enabled schedules and extract names\n"
    "iics schedule list -o json \\\n"
    "  | jq -r '.[] | select(.status==\"ENABLED\") | .name'\n\n"
    "# Count objects in a project\n"
    "iics objects list --project Prod -o json | jq length"
), Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.1), font_size=11)

add_text_box(s,
    "stdout is always clean data.  All errors and diagnostics go to stderr.",
    Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
    font_size=14, italic=True, color=MID_BLUE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – Secure Credentials
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Secure Credential Management", "Credentials never appear on the command line")
footer(s)

add_text_box(s, "Official CLI – credentials on every command line",
             Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=RED)
code_block(s, (
    "# Visible in ps aux, shell history, CI logs\n"
    "iics export \\\n"
    "  -u admin@company.com \\\n"
    "  -p MyP@ssword123 \\\n"
    "  -r us \\\n"
    "  -a \"Prod/CriticalFlow.Process\""
), Inches(0.5), Inches(1.65), Inches(5.8), Inches(1.8), font_size=11)

add_text_box(s, "iics-cli – config file + env var + OS keychain",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
code_block(s, (
    "# ~/.iics/config.yaml  (chmod 600)\n"
    "defaultProfile: prod\n"
    "profiles:\n"
    "  prod:\n"
    "    region: USE4\n"
    "    username: admin@company.com\n"
    "    password: \"@keyring\"  # stored in OS keychain\n\n"
    "# Password via env var – never in a file\n"
    "export IICS_PASSWORD='MyP@ssword123'\n"
    "iics export create --name nightly"
), Inches(7.0), Inches(1.65), Inches(5.8), Inches(2.2), font_size=11)

bullet_list(s, [
    "Config file protected with OS file permissions (chmod 600)",
    "Password stored in macOS Keychain or Linux Secret Service (GNOME Keyring / KWallet)",
    "Password sourced from IICS_PASSWORD env var (highest precedence)",
    "Sentinel value @keyring in config – real secret never written to disk",
    "Compatible with HashiCorp Vault, AWS Secrets Manager, GitHub Secrets",
], Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.6), font_size=14)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – Named Profiles + Interactive Wizard
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Named Environment Profiles", "Interactive wizard or manual config – switch with a single flag")
footer(s)

add_text_box(s, "Interactive profile setup",
             Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
code_block(s, (
    "# First-run wizard – guided setup\n"
    "iics profile add\n\n"
    "# Or add a named profile\n"
    "iics profile add prod\n\n"
    "# List / edit / delete profiles\n"
    "iics profile list\n"
    "iics profile edit dev\n"
    "iics profile set-default prod"
), Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.3), font_size=11)

add_text_box(s, "~/.iics/config.yaml",
             Inches(0.5), Inches(4.05), Inches(5.8), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
code_block(s, (
    "defaultProfile: dev\n"
    "profiles:\n"
    "  dev:\n"
    "    region: USE4\n"
    "    username: dev-user@company.com\n"
    "  prod:\n"
    "    region: USE4\n"
    "    username: prod-admin@company.com\n"
    "    password: \"@keyring\"\n"
    "  emea:\n"
    "    region: EMEA\n"
    "    username: emea-admin@company.com"
), Inches(0.5), Inches(4.45), Inches(5.8), Inches(2.5), font_size=11)

add_text_box(s, "Switching environments",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
code_block(s, (
    "# Use the --profile flag\n"
    "iics user list --profile dev\n"
    "iics user list --profile prod\n\n"
    "# Or the IICS_PROFILE env var\n"
    "IICS_PROFILE=emea iics export create \\\n"
    "  --name weekly-backup\n\n"
    "# Perfect for deployment scripts\n"
    "for env in dev staging prod; do\n"
    "  iics export create --profile $env \\\n"
    "    --name \"snapshot-$env\"\n"
    "done"
), Inches(7.0), Inches(1.6), Inches(5.8), Inches(3.0), font_size=11)

bullet_list(s, [
    "Separate credentials per environment – no copy-paste mistakes",
    "Profile overridden by IICS_PROFILE env var for CI/CD",
    "defaultProfile used when no --profile flag is given",
    "iics profile add wizard prompts for region, username, and keychain storage",
], Inches(7.0), Inches(4.75), Inches(5.8), Inches(2.0), font_size=13)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – Session Caching
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Session Caching", "Authenticate once, reuse for 30 minutes")
footer(s)

add_text_box(s, "Official CLI – authenticates on every command",
             Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=RED)
bullet_list(s, [
    "Login round-trip added to every command",
    "Risk of account lockout during scripted runs",
    "Slower in pipelines with many sequential commands",
], Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.0), font_size=14, color=RED)

add_text_box(s, "iics-cli – session token cached with 30-min TTL",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
code_block(s, (
    "iics login                   # authenticates once\n"
    "iics user list               # reuses session\n"
    "iics role list               # reuses session\n"
    "iics connection list         # reuses session\n"
    "iics export create --name x  # reuses session\n\n"
    "# Separate cached session per profile\n"
    "iics user list --profile prod  # own token"
), Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.8), font_size=11)

bullet_list(s, [
    "Sessions cached in ~/.iics/sessions.yaml",
    "Each profile has its own independent session entry",
    "Expired or invalid sessions trigger automatic re-login transparently",
    "401 responses auto-retry once with a fresh login – no manual intervention",
], Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.5), font_size=15)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – API Coverage
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Broad API Coverage", "27+ resource types vs. 5 operations in the official tool")
footer(s)

add_table(s,
    ["Resource", "Operations"],
    [
        ["objects",       "list, dependencies"],
        ["lookup",        "resolve IDs, names, paths"],
        ["connection",    "list, get, create, update, delete"],
        ["export",        "run, create, start, status, download"],
        ["import",        "run, upload, start, status, download-log"],
        ["publish",       "run, start, status"],
        ["unpublish",     "run, start, status"],
        ["package",       "expand, create, dependencies"],
        ["release",       "manifest, validate, plan"],
        ["schedule",      "list, get, create, update, delete"],
        ["project",       "create, update, delete"],
        ["folder",        "create, update, delete"],
        ["profile",       "add, edit, list, delete, set-default, show"],
    ],
    Inches(0.4), Inches(1.25), Inches(6.0), Inches(5.7))

add_table(s,
    ["Resource", "Operations"],
    [
        ["user",          "list, get, create, update, delete"],
        ["usergroup",     "list, get, create, update, delete"],
        ["role",          "list, get, create, update, delete"],
        ["privilege",     "list"],
        ["runtime",       "list, get, create, update"],
        ["agent",         "list, get, details, start, stop"],
        ["tag",           "assign, remove"],
        ["permission",    "get, set, delete"],
        ["activitylog",   "list, get"],
        ["auditlog",      "list"],
        ["securitylog",   "list"],
        ["metering",      "get, download"],
        ["sourcecontrol", "checkout, checkin, pull, commit"],
        ["state",         "fetch, load"],
    ],
    Inches(6.9), Inches(1.25), Inches(6.0), Inches(5.7))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – Package Command
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Package Command", "Expand, create, and inspect IICS export packages locally")
footer(s)

add_text_box(s, "Local file operations – no API calls required",
             Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.4),
             font_size=15, bold=True, color=DARK_BLUE)

add_text_box(s, "Expand, Create, Inspect",
             Inches(0.5), Inches(1.7), Inches(5.8), Inches(0.35),
             font_size=13, bold=True, color=MID_BLUE)
code_block(s, (
    "# Expand a ZIP into a source-control-friendly tree\n"
    "iics package expand --source export.zip \\\n"
    "  --target ./workspace/Prod\n\n"
    "# Create a ZIP from an expanded workspace\n"
    "iics package create --source ./workspace/Prod \\\n"
    "  --target deploy.zip\n\n"
    "# Selective packaging via manifest\n"
    "iics package create --source ./workspace/Prod \\\n"
    "  --manifest-file manifest.csv \\\n"
    "  --target deploy.zip\n\n"
    "# Inspect dependencies in a package\n"
    "iics package dependencies --source ./workspace/Prod \\\n"
    "  -o json | jq '.[].location'"
), Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5), font_size=11)

add_text_box(s, "Key capabilities",
             Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
             font_size=13, bold=True, color=MID_BLUE)
bullet_list(s, [
    "Expand: unzip to Project/Folder/Asset tree for git version control",
    "Create: re-assemble a deployment ZIP from workspace directory",
    "Selective create: drive asset selection via CSV/JSON/YAML/TXT manifest",
    "Manifest auto-detects format – stdin or --manifest-file (-m)",
    "Transitive dependency resolution – auto-include connections and agents",
    "Generates ContentsofExportPackage_*.csv and exportMetadata.v2.json",
    "Computes xexportPackage.chksum checksum file",
    "Strip CurrentServerDateTime to enable clean git diffs",
    "Backward-compatible with official CLI asset list text files",
], Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.5), font_size=13)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – Release Workflow
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Release Workflow", "CI/CD manifests and per-environment deployment plans")
footer(s)

add_text_box(s, "Three-step release pipeline",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.35),
             font_size=15, bold=True, color=DARK_BLUE)

code_block(s, (
    "# 1. Generate a release manifest from dependency graph\n"
    "iics release manifest \\\n"
    "  --project Prod \\\n"
    "  --output release-manifest.yaml\n\n"
    "# 2. Validate the manifest against each target\n"
    "iics release validate \\\n"
    "  --manifest release-manifest.yaml \\\n"
    "  --profile staging\n\n"
    "# 3. Generate per-environment deployment plans\n"
    "iics release plan \\\n"
    "  --manifest release-manifest.yaml \\\n"
    "  --targets dev,staging,prod \\\n"
    "  --output-dir ./plans"
), Inches(0.5), Inches(1.65), Inches(6.0), Inches(4.5), font_size=11)

add_text_box(s, "CI/CD integration",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.35),
             font_size=14, bold=True, color=MID_BLUE)
bullet_list(s, [
    "release manifest – scans object dependencies and outputs a YAML manifest",
    "release validate – checks each asset exists in the target org",
    "release plan – generates one plan CSV per deployment target",
    "IICS_VALID_DEPLOY_TARGETS – allowlist of approved targets",
    "IICS_TARGET_PROFILE_MAP – maps target names to profiles",
    "Per-target CI credentials: IICS_USER_<TARGET>, IICS_PWD_<TARGET>",
    "Works with GitHub Actions, Jenkins, GitLab CI, and any shell pipeline",
    "Publish order enforcement – dependencies deployed before dependents",
], Inches(7.0), Inches(1.65), Inches(5.8), Inches(4.5), font_size=13)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – Publish & Unpublish
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Publish & Unpublish CAI Assets", "Control asset lifecycle on the runtime")
footer(s)

add_text_box(s, "Publish assets to the runtime",
             Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.35),
             font_size=14, bold=True, color=GREEN)
code_block(s, (
    "# Publish immediately and wait for completion\n"
    "iics publish run \\\n"
    "  --ids <asset-id-1>,<asset-id-2>\n\n"
    "# Start async publish job\n"
    "iics publish start --ids <asset-id-1>\n\n"
    "# Check publish status\n"
    "iics publish status --id <job-id>"
), Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.4), font_size=11)

add_text_box(s, "Unpublish assets from the runtime",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.35),
             font_size=14, bold=True, color=ORANGE)
code_block(s, (
    "# Unpublish immediately and wait for completion\n"
    "iics unpublish run \\\n"
    "  --ids <asset-id-1>,<asset-id-2>\n\n"
    "# Start async unpublish job\n"
    "iics unpublish start --ids <asset-id-1>\n\n"
    "# Check unpublish status\n"
    "iics unpublish status --id <job-id>"
), Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.4), font_size=11)

add_text_box(s, "Supported asset types",
             Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
add_table(s,
    ["Asset Type", "Publishable", "Notes"],
    [
        ["PROCESS",              "Yes", "Application Integration processes"],
        ["GUIDE",                "Yes", "Application Integration guides"],
        ["TASKFLOW",             "Yes", "Taskflows"],
        ["AI_CONNECTION",        "Yes", "Application Integration connections"],
        ["AI_SERVICE_CONNECTOR", "Yes", "Application Integration service connectors"],
        ["DTEMPLATE",            "No",  "Mappings – deployed via import, not publish"],
    ],
    Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.3))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 – Table Themes
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Table Themes & Output Styles", "Human-readable and automation-friendly in the same tool")
footer(s)

add_text_box(s, "Six built-in table themes – global or per-command",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
             font_size=15, bold=True, color=DARK_BLUE)

add_table(s,
    ["Theme", "Description", "Color", "Always Renders"],
    [
        ["default",  "Unicode rounded borders, cyan bold headers",        "Yes (TTY)", "No"],
        ["minimal",  "No borders, colored bold headers, unicode underline","Yes (TTY)", "No"],
        ["compact",  "No borders, gray bold headers, 1-space column gap", "Yes (TTY)", "No"],
        ["plain",    "ASCII borders, no color – auto-used for non-TTY",   "No",        "No"],
        ["markdown", "GitHub-flavored markdown table",                    "No",        "Yes"],
        ["gh",       "GitHub CLI-style: no borders, plain headers",       "No",        "Yes"],
    ],
    Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.8))

add_text_box(s, "Configure theme globally or per invocation",
             Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.35),
             font_size=14, bold=True, color=DARK_BLUE)
code_block(s, (
    "# Set a default theme in config\n"
    "# ~/.iics/config.yaml\n"
    "style:\n"
    "  theme: minimal\n"
    "  noColor: false\n"
    "  headerColor: \"6\"  # lipgloss: 6=cyan, 244=gray, #FF0000=hex\n\n"
    "# Override per command\n"
    "iics user list --theme markdown   # paste into GitHub / Confluence\n"
    "iics connection list --theme gh   # GitHub CLI style\n"
    "IICS_THEME=compact iics role list # env var override"
), Inches(0.5), Inches(5.05), Inches(12.3), Inches(2.0), font_size=11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 – Region Support
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Comprehensive Region Support", "All IICS pod regions – not just us / eu / ap")
footer(s)

add_text_box(s, "Official CLI – 3 hardcoded regions",
             Inches(0.5), Inches(1.2), Inches(4.0), Inches(0.4),
             font_size=14, bold=True, color=RED)
bullet_list(s, ["-r us", "-r eu", "-r ap"],
            Inches(0.5), Inches(1.65), Inches(4.0), Inches(1.2),
            font_size=14, color=RED)

add_text_box(s, "iics-cli – full pod registry",
             Inches(5.0), Inches(1.2), Inches(7.8), Inches(0.4),
             font_size=14, bold=True, color=GREEN)

pods = [
    ["US, USW1, USE2, USW3, USE4, USW5, USE6", "US (various pods)"],
    ["USW1-1, USW3-1",                          "US West (dm1 hosts)"],
    ["USW1-2",                                  "US West (dm2 host)"],
    ["CAC1",                                    "Canada"],
    ["EMEA, EMWE1",                             "Europe / Middle East / Africa"],
    ["APSE1, APJ",                              "Asia Pacific"],
    ["APNE1",                                   "Asia Pacific North-East"],
]
add_table(s,
    ["Region Code(s)", "Geography"],
    pods,
    Inches(5.0), Inches(1.65), Inches(7.8), Inches(3.1))

add_text_box(s,
    "loginUrl and baseApiUrl auto-populated on first iics login – no manual URL entry needed.",
    Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
    font_size=14, italic=True, color=MID_BLUE)

code_block(s, (
    "# Use a built-in region\n"
    "iics user list --profile emea   # region: EMEA in config\n\n"
    "# Override with explicit login URL (future pods / custom deployments)\n"
    "# loginUrl: https://dm-us.informaticacloud.com/saas/public/core/v3/login"
), Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.4), font_size=11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 – Modern Design & Broad Platform Support
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Modern Design & Broad Platform Support", "Built for today's DevOps toolchains")
footer(s)

add_text_box(s, "Official CLI – Limited and outdated",
             Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=RED)
bullet_list(s, [
    "No longer maintained – frozen in time",
    "Covers only export/import/publish operations",
    "Limited region support (us/eu/ap only)",
    "Complex credential and profile management",
    "No modern automation-friendly features",
], Inches(0.5), Inches(1.7), Inches(5.8), Inches(3.5), font_size=15, color=RED)

add_text_box(s, "iics-cli – Modern, actively maintained",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
bullet_list(s, [
    "Actively maintained with regular updates",
    "27+ resource types with full CRUD operations",
    "All IICS pod regions supported",
    "Named profiles, OS keychain, and session caching",
    "~8 MB statically-linked Go binary – no runtime dependencies",
], Inches(7.0), Inches(1.7), Inches(5.8), Inches(3.5), font_size=15, color=GREEN)

code_block(s, (
    "# Install – one line\n"
    "curl -L https://github.com/jbrazda/iics-cli/releases/latest/download/iics_linux_amd64 \\\n"
    "  -o /usr/local/bin/iics && chmod +x /usr/local/bin/iics\n\n"
    "# Docker – copy single binary\n"
    "COPY iics /usr/local/bin/iics"
), Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), font_size=11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 – Structured Commands + Help
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Structured Subcommands & Built-in Help", "kubectl / gh / aws style resource-verb hierarchy")
footer(s)

code_block(s, (
    "# Resource  Verb  Flags\n"
    "iics user   list\n"
    "iics user   get   --id abc123\n"
    "iics user   create --from-file user.json\n\n"
    "iics connection list\n"
    "iics connection update --id xyz --from-file conn.json\n\n"
    "iics export run --name nightly --project Prod\n"
    "iics export status --id 12345\n"
    "iics export download --id 12345 --output export.zip\n\n"
    "iics package expand --source export.zip --target ./workspace\n"
    "iics package create --source ./workspace --manifest-file plan.csv\n\n"
    "iics release manifest --project Prod --output manifest.yaml\n"
    "iics release plan --manifest manifest.yaml --targets dev,prod\n\n"
    "iics sourcecontrol checkin --project Prod --comment \"release 2.4\"\n"
    "iics state fetch --project Prod --output state.json"
), Inches(0.5), Inches(1.3), Inches(7.5), Inches(5.5), font_size=10)

add_text_box(s, "Built-in help at every level",
             Inches(8.2), Inches(1.3), Inches(4.8), Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE)
code_block(s, (
    "iics --help\n\n"
    "iics user --help\n\n"
    "iics export create --help\n\n"
    "iics package --help"
), Inches(8.2), Inches(1.75), Inches(4.8), Inches(1.7), font_size=11)

bullet_list(s, [
    "Same pattern as kubectl, gh, aws – familiar to DevOps teams",
    "Tab completion: bash, zsh, fish, powershell",
    "Per-command docs in docs/documentation/",
    "Full flag reference in built-in --help",
    "--debug prints full HTTP trace to stderr",
], Inches(8.2), Inches(3.6), Inches(4.8), Inches(2.8), font_size=13)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 – Error Reporting
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Better Error Reporting", "Structured API errors with HTTP context – all on stderr")
footer(s)

add_text_box(s, "Official CLI – vague error messages",
             Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=RED)
code_block(s, (
    "$ iics export -u user -p pass -r us \\\n"
    "  -a \"Project/Missing.Process\"\n\n"
    "Error occurred.\n"
    "exit code: 0"
), Inches(0.5), Inches(1.65), Inches(5.8), Inches(1.5), font_size=11)

add_text_box(s, "iics-cli – structured errors to stderr",
             Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
code_block(s, (
    "$ iics export create --name bad --project X\n\n"
    "Error: API error 400: project 'X' not found\n"
    "  HTTP 400 Bad Request\n"
    "  {\n"
    "    \"error\": {\n"
    "      \"code\": \"ICS-012345\",\n"
    "      \"message\": \"project 'X' not found\"\n"
    "    }\n"
    "  }\n"
    "exit code: 1"
), Inches(7.0), Inches(1.65), Inches(5.8), Inches(2.5), font_size=11)

bullet_list(s, [
    "Error summary always on stderr – stdout stays clean for piping",
    "--debug flag prints full request body for deep troubleshooting",
    "--verbose flag enables extra HTTP detail",
    "API error code, HTTP status, and response body all included",
    "Usage errors (exit 2) print a hint without cluttering stdout",
], Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.5), font_size=15)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 – Feature Comparison Summary
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Feature Comparison Summary")
footer(s)

add_table(s,
    ["Capability", "Official Informatica CLI", "iics-cli"],
    [
        ["Actively maintained",               "No",      "Yes"],
        ["POSIX exit codes",                  "No",      "Yes"],
        ["JSON / CSV / YAML output",          "No",      "Yes"],
        ["Pipe-friendly stdout/stderr split", "No",      "Yes"],
        ["Credentials in config file",        "No",      "Yes"],
        ["OS Keychain integration",           "No",      "Yes"],
        ["Named environment profiles",        "No",      "Yes"],
        ["Interactive profile wizard",        "No",      "Yes"],
        ["Session caching",                   "No",      "Yes"],
        ["Package expand / create",           "No",      "Yes"],
        ["Selective packaging via manifest",  "No",      "Yes"],
        ["CI/CD release manifest & plan",     "No",      "Yes"],
        ["Publish / Unpublish",               "Partial", "Yes"],
        ["User & group management",           "No",      "Yes"],
        ["Connection CRUD",                   "No",      "Yes"],
        ["Schedule CRUD",                     "No",      "Yes"],
        ["Agent & runtime management",        "No",      "Yes"],
        ["Activity / Audit / Security logs",  "No",      "Yes"],
        ["Source control operations",         "No",      "Yes"],
        ["All IICS regions",                  "No",      "Yes"],
        ["Table themes & styles",             "No",      "Yes"],
        ["Built-in contextual help",          "Minimal", "Yes"],
    ],
    Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.9))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 19 – Getting Started
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Getting Started", "Install, configure, and run your first command in minutes")
footer(s)

code_block(s, (
    "# 1. Download the binary for your platform\n"
    "curl -L https://github.com/jbrazda/iics-cli/releases/latest/download/iics_linux_amd64 \\\n"
    "  -o /usr/local/bin/iics && chmod +x /usr/local/bin/iics\n\n"
    "# 2. Create a profile interactively (guided wizard)\n"
    "iics profile add\n\n"
    "#    -- OR create ~/.iics/config.yaml manually --\n"
    "mkdir -p ~/.iics && chmod 700 ~/.iics\n"
    "cat > ~/.iics/config.yaml << 'EOF'\n"
    "defaultProfile: default\n"
    "profiles:\n"
    "  default:\n"
    "    region: USE4\n"
    "    username: your-email@company.com\n"
    "EOF\n"
    "chmod 600 ~/.iics/config.yaml\n\n"
    "# 3. Set your password in the environment (or use the keychain via wizard)\n"
    "export IICS_PASSWORD='your-password'\n\n"
    "# 4. Login (optional – cached for 30 min)\n"
    "iics login\n\n"
    "# 5. Start exploring\n"
    "iics user list\n"
    "iics connection list -o json | jq '.[].name'\n"
    "iics objects list --type MTT --output json\n"
    "iics export create --name my-first-export --project MyProject"
), Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), font_size=11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 20 – Thank You / Questions
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
add_rect(s, 0, Inches(3.5), SLIDE_W, Inches(0.06), ORANGE)

add_text_box(s, "Questions?",
             Inches(1), Inches(1.5), Inches(11), Inches(1.0),
             font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s, "github.com/jbrazda/iics-cli",
             Inches(1), Inches(2.8), Inches(11), Inches(0.6),
             font_size=22, color=RGBColor(0xFF, 0xAA, 0x44), align=PP_ALIGN.CENTER)

add_text_box(s, "Source code  |  Releases  |  Issues  |  Documentation",
             Inches(1), Inches(3.6), Inches(11), Inches(0.5),
             font_size=16, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

add_text_box(s, "docs/documentation/  in the repository for full per-command reference",
             Inches(1), Inches(4.3), Inches(11), Inches(0.5),
             font_size=14, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = "iics-cli-presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
